from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ============================================================
# カラーパレット
# ============================================================
COLOR_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 紺黒（背景）
COLOR_ACCENT    = RGBColor(0xE8, 0x3A, 0x3A)   # 赤（危機感）
COLOR_ACCENT2   = RGBColor(0xF5, 0xA6, 0x23)   # オレンジ（強調）
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)   # 薄いグレー文字
COLOR_CARD_BG   = RGBColor(0x2A, 0x2A, 0x45)   # カード背景
COLOR_TABLE_HDR = RGBColor(0xE8, 0x3A, 0x3A)
COLOR_TABLE_ROW = RGBColor(0x22, 0x22, 0x3A)
COLOR_TABLE_ROW2= RGBColor(0x2E, 0x2E, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 完全ブランク


# ============================================================
# ユーティリティ
# ============================================================
def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(14), bold=False, color=COLOR_WHITE,
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_text_lines(slide, lines, l, t, w, h,
                   font_size=Pt(13), bold=False, color=COLOR_WHITE,
                   align=PP_ALIGN.LEFT, line_spacing=Pt(6)):
    """複数行テキスト（リスト）を段落として追加"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


# ============================================================
# SLIDE 1
# ============================================================
slide1 = prs.slides.add_slide(BLANK)

# 背景
add_rect(slide1, 0, 0, 13.33, 7.5, fill_color=COLOR_BG_DARK)

# 左側の赤い縦ライン
add_rect(slide1, 0.35, 0.4, 0.06, 6.7, fill_color=COLOR_ACCENT)

# スライド番号バッジ
badge = add_rect(slide1, 0.5, 0.38, 0.9, 0.38, fill_color=COLOR_ACCENT)
add_text(slide1, "01 / 02", 0.5, 0.38, 0.9, 0.38,
         font_size=Pt(9), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# メインタイトル
add_text(slide1,
         "毎月リセット。この営業、いつまで続けられるか。",
         0.55, 0.85, 12.3, 0.65,
         font_size=Pt(26), bold=True, color=COLOR_WHITE)

# 赤い下線
add_rect(slide1, 0.55, 1.52, 12.3, 0.04, fill_color=COLOR_ACCENT)

# ── ボディ左カラム ──────────────────────────────────────────
# リード文
add_text_lines(slide1,
    ["正直に言う。",
     "",
     "今のBP事業部の売上は、「たまたま」で成り立っている。",
     "大型案件が取れた月は達成。取れなかった月は未達。",
     "翌月の見込みは、月が変わるまでわからない。",
     "",
     "これは「営業している」のではなく、「注文を待っている」だ。"],
    0.55, 1.65, 5.8, 2.0,
    font_size=Pt(12.5), color=COLOR_LIGHT)

# 「注文を待っている」部分を強調（別テキストで上書き）
add_text(slide1, "「注文を待っている」",
         2.92, 2.48, 3.2, 0.38,
         font_size=Pt(12.5), bold=True, color=COLOR_ACCENT2)

# 数字カード
card_t = 3.75
card_data = [
    ("月次目標", "3,000万円", ""),
    ("現状見込み", "800〜900万円", ""),
    ("毎月の祈り", "▲2,000万円", "どこからか来ることを"),
]
card_colors = [COLOR_CARD_BG, COLOR_CARD_BG, RGBColor(0x3A, 0x18, 0x18)]
card_w = 1.72
for i, (label, val, sub) in enumerate(card_data):
    cx = 0.55 + i * (card_w + 0.12)
    add_rect(slide1, cx, card_t, card_w, 1.15,
             fill_color=card_colors[i],
             line_color=COLOR_ACCENT if i == 2 else None,
             line_width=Pt(1.5) if i == 2 else Pt(0))
    add_text(slide1, label, cx+0.1, card_t+0.07, card_w-0.2, 0.28,
             font_size=Pt(9), color=COLOR_LIGHT)
    add_text(slide1, val, cx+0.1, card_t+0.33, card_w-0.2, 0.45,
             font_size=Pt(18), bold=True,
             color=COLOR_ACCENT if i == 2 else COLOR_WHITE)
    if sub:
        add_text(slide1, sub, cx+0.1, card_t+0.80, card_w-0.2, 0.28,
                 font_size=Pt(8.5), color=COLOR_LIGHT, italic=True)

# ── ボディ右カラム ──────────────────────────────────────────
# 区切り縦線
add_rect(slide1, 6.65, 1.65, 0.03, 4.8, fill_color=RGBColor(0x44, 0x44, 0x66))

# 右側リード
add_text(slide1, "もっと深刻な事実がある。",
         6.8, 1.65, 6.2, 0.35,
         font_size=Pt(13), bold=True, color=COLOR_ACCENT2)

add_text_lines(slide1,
    ["直近1年でアクティブな取引先は約 290社。",
     "過去10年、数百社の法人と関係を築いてきた。",
     "",
     "なのに——"],
    6.8, 2.05, 6.2, 1.0,
    font_size=Pt(12), color=COLOR_LIGHT)

# 問いかけブロック
questions = [
    "その会社が年間いくら発注しているか、把握できているか？",
    "他の部署に発注担当者がいないか、確認したことがあるか？",
    "なぜ先月より注文が減ったのか、直接聞いたことがあるか？",
]
q_t = 3.18
for q in questions:
    add_rect(slide1, 6.8, q_t, 6.2, 0.36,
             fill_color=RGBColor(0x28, 0x28, 0x42))
    add_text(slide1, "▶  " + q, 6.88, q_t+0.04, 6.05, 0.32,
             font_size=Pt(11), color=COLOR_LIGHT)
    q_t += 0.42

add_text(slide1, "ほとんど、できていない。",
         6.8, q_t + 0.1, 6.2, 0.38,
         font_size=Pt(14), bold=True, color=COLOR_ACCENT)

# 結論バナー
add_rect(slide1, 0, 6.72, 13.33, 0.78, fill_color=COLOR_ACCENT)
add_text(slide1,
         "このままでは売上は上がらない。チームも守れない。営業チーム維持には月5,000万円規模の売上が必要だ。今のやり方で、それが実現できるか——答えは全員わかっているはずだ。",
         0.4, 6.77, 12.53, 0.65,
         font_size=Pt(11.5), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2
# ============================================================
slide2 = prs.slides.add_slide(BLANK)

# 背景
add_rect(slide2, 0, 0, 13.33, 7.5, fill_color=COLOR_BG_DARK)

# 左の赤縦ライン
add_rect(slide2, 0.35, 0.4, 0.06, 6.7, fill_color=COLOR_ACCENT)

# スライド番号バッジ
add_rect(slide2, 0.5, 0.38, 0.9, 0.38, fill_color=COLOR_ACCENT)
add_text(slide2, "02 / 02", 0.5, 0.38, 0.9, 0.38,
         font_size=Pt(9), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# タイトル
add_text(slide2,
         "「注文を待つ営業」を終わりにする。",
         0.55, 0.85, 12.3, 0.65,
         font_size=Pt(26), bold=True, color=COLOR_WHITE)

add_rect(slide2, 0.55, 1.52, 12.3, 0.04, fill_color=COLOR_ACCENT)

# サブタイトル
add_text(slide2,
         "新しい戦略は一言で言える。既存顧客を、会社ごと握る。 ——これが「アカウント営業」だ。",
         0.55, 1.62, 12.3, 0.42,
         font_size=Pt(13), bold=False, color=COLOR_ACCENT2)

# ── 左：Before / After ────────────────────────────────────
add_text(slide2, "何が変わるのか", 0.55, 2.15, 5.8, 0.32,
         font_size=Pt(11), bold=True, color=COLOR_LIGHT)

rows_ba = [
    ("これまで",          "これから"),
    ("担当者一人から注文をもらう",   "会社全体の発注を継続的に取り続ける"),
    ("単発スポットを積み上げる",     "ストック型の安定売上を構築する"),
    ("注文が来るのを待つ",          "ニーズを掘り起こし、提案しにいく"),
    ("案件単位で管理",              "企業単位のアカウントで管理"),
]
row_h = 0.40
tbl_t = 2.50
for ri, (before, after) in enumerate(rows_ba):
    is_header = ri == 0
    bg_b = COLOR_TABLE_HDR if is_header else (COLOR_TABLE_ROW if ri % 2 == 1 else COLOR_TABLE_ROW2)
    bg_a = COLOR_TABLE_HDR if is_header else (COLOR_TABLE_ROW if ri % 2 == 1 else COLOR_TABLE_ROW2)
    fc   = COLOR_WHITE
    fs   = Pt(10.5) if not is_header else Pt(11)
    bold = is_header

    add_rect(slide2, 0.55, tbl_t + ri*row_h, 2.75, row_h, fill_color=bg_b)
    add_text(slide2, before, 0.65, tbl_t + ri*row_h + 0.05, 2.55, row_h-0.1,
             font_size=fs, bold=bold, color=fc)

    # 矢印
    if not is_header:
        add_text(slide2, "→", 3.32, tbl_t + ri*row_h + 0.05, 0.3, row_h-0.1,
                 font_size=Pt(14), bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide2, 3.62, tbl_t + ri*row_h, 2.75, row_h, fill_color=bg_a)
    add_text(slide2, after, 3.72, tbl_t + ri*row_h + 0.05, 2.55, row_h-0.1,
             font_size=fs, bold=bold, color=fc)

# ── 右：セグメント別戦術 ───────────────────────────────────
add_rect(slide2, 6.65, 2.15, 0.03, 4.3, fill_color=RGBColor(0x44, 0x44, 0x66))

add_text(slide2, "顧客タイプ別のアプローチ", 6.82, 2.15, 6.1, 0.32,
         font_size=Pt(11), bold=True, color=COLOR_LIGHT)

seg_data = [
    {
        "tag":   "集中型",
        "desc":  "特定の担当者1人に発注が集中",
        "goal":  "目指すゴール：年間発注を握るパートナー化",
        "tactic":["・年間見積を提示し、定例会を設定する",
                  "・付加価値提供で競合を排除する",
                  "・キーパーソンとの関係を組織単位に広げる"],
        "color": COLOR_ACCENT,
    },
    {
        "tag":   "分散型",
        "desc":  "社内の多くの社員がバラバラに発注",
        "goal":  "目指すゴール：社内インフラ化（発注一本化）",
        "tactic":["・利用状況をヒアリングし、課題を把握する",
                  "・発注集約＋コスト削減の提案を持ち込む",
                  "・利用ガイド作成など、社内展開を支援する"],
        "color": COLOR_ACCENT2,
    },
]
card_t2 = 2.55
for seg in seg_data:
    add_rect(slide2, 6.82, card_t2, 6.1, 1.75,
             fill_color=COLOR_CARD_BG,
             line_color=seg["color"], line_width=Pt(1.5))
    # タグ
    add_rect(slide2, 6.82, card_t2, 1.1, 0.30, fill_color=seg["color"])
    add_text(slide2, seg["tag"], 6.82, card_t2, 1.1, 0.30,
             font_size=Pt(10), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide2, seg["desc"], 7.98, card_t2+0.02, 4.8, 0.28,
             font_size=Pt(10), color=COLOR_LIGHT)
    add_text(slide2, seg["goal"], 6.92, card_t2+0.34, 5.8, 0.28,
             font_size=Pt(10.5), bold=True, color=seg["color"])
    add_text_lines(slide2, seg["tactic"],
                   6.92, card_t2+0.62, 5.8, 1.0,
                   font_size=Pt(10), color=COLOR_LIGHT)
    card_t2 += 1.92

# ── 下部：実行メッセージ ──────────────────────────────────
add_rect(slide2, 0.55, 6.3, 12.45, 0.04, fill_color=RGBColor(0x44, 0x44, 0x66))

add_text_lines(slide2,
    ["やることは決まった。対象リストも出す。スクリプトも作る。あとは動くだけだ。",
     "下期終了時に「あのとき変わってよかった」と言えるかどうかは、今週の行動で決まる。"],
    0.55, 6.42, 12.45, 0.9,
    font_size=Pt(11.5), bold=False, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
out_path = "/Volumes/SSD-PGU3C/営業推進/営業戦略スライド/BP事業部_営業戦略_大前提ストーリー.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
