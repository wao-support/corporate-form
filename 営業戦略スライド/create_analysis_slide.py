from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# カラーパレット
# ============================================================
COLOR_BG        = RGBColor(0x10, 0x10, 0x20)
COLOR_ACCENT    = RGBColor(0xE8, 0x3A, 0x3A)
COLOR_ACCENT2   = RGBColor(0xF5, 0xA6, 0x23)
COLOR_GREEN     = RGBColor(0x2E, 0xCC, 0x71)
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xBB, 0xBB, 0xCC)
COLOR_MUTED     = RGBColor(0x77, 0x77, 0x99)
COLOR_CARD      = RGBColor(0x1E, 0x1E, 0x35)
COLOR_CARD2     = RGBColor(0x28, 0x28, 0x45)
COLOR_DIVIDER   = RGBColor(0x44, 0x44, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================
# ユーティリティ
# ============================================================
def rect(l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = lw
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def txt(text, l, t, w, h, size=Pt(13), bold=False, color=COLOR_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def txt_lines(lines, l, t, w, h, size=Pt(12), bold=False,
              color=COLOR_WHITE, align=PP_ALIGN.LEFT, spacing=Pt(4)):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = spacing
        r = p.add_run()
        r.text = line
        r.font.size = size
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


# ============================================================
# 背景
# ============================================================
rect(0, 0, 13.33, 7.5, fill=COLOR_BG)

# 上部アクセントライン
rect(0, 0, 13.33, 0.06, fill=COLOR_ACCENT)

# 左縦ライン
rect(0.32, 0.18, 0.055, 7.1, fill=COLOR_ACCENT)


# ============================================================
# ヘッダー
# ============================================================
txt("現状分析", 0.5, 0.14, 2.5, 0.32,
    size=Pt(9), bold=True, color=COLOR_ACCENT)

txt("営業は「売上を作っている」のか？　データで検証した。",
    0.5, 0.44, 10.0, 0.50,
    size=Pt(20), bold=True, color=COLOR_WHITE)

rect(0.5, 0.96, 12.5, 0.03, fill=COLOR_DIVIDER)

# ソースタグ
rect(11.2, 0.18, 1.95, 0.28, fill=RGBColor(0x2A,0x2A,0x45))
txt("出典：Salesforce活動ログ分析", 11.25, 0.20, 1.85, 0.26,
    size=Pt(8), color=COLOR_MUTED)


# ============================================================
# 中央：衝撃の一文
# ============================================================
rect(1.5, 1.08, 10.33, 1.22,
     fill=RGBColor(0x2A, 0x08, 0x08),
     line=COLOR_ACCENT, lw=Pt(2))

txt("「営業が売上を作っているのか」を確かめに行ったら、",
    1.65, 1.13, 10.0, 0.38,
    size=Pt(13), color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

txt("「営業は、売上を処理していた」",
    1.65, 1.50, 10.0, 0.55,
    size=Pt(22), bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

rect(2.8, 1.47, 7.73, 0.03, fill=RGBColor(0x66,0x22,0x22))


# ============================================================
# 3カラム：Evidence
# ============================================================
col_titles = ["① 初回活動の実態", "② 営業手法の分布", "③ 売上の発生源"]
col_icons  = ["📥", "📧", "🌐"]
col_bodies = [
    ["多くの初回活動が「Re:メール」", "＝ 問い合わせへの返信", "",
     "▶  問い合わせが先にあり", "　  営業はそれに対応している"],
    ["活動の大半がメール対応", "電話（コール）は極めて少ない", "",
     "▶  アウトバウンド営業は", "　  ほぼ行われていない"],
    ["活動ログなし（#N/A）の", "受注が一定数存在", "",
     "▶  Web自然流入・紹介経由の", "　  案件が相当数を占める"],
]
col_conclusions = [
    "営業起点ではなく受動対応",
    "能動的な案件創出が弱い",
    "営業が直接創出した売上は不明",
]
col_w = 3.8
col_gap = 0.28
col_start = 0.5

for i in range(3):
    cx = col_start + i * (col_w + col_gap)
    ct = 2.42

    # カード背景
    rect(cx, ct, col_w, 2.88,
         fill=COLOR_CARD2,
         line=COLOR_DIVIDER, lw=Pt(0.75))

    # カラムタイトル
    rect(cx, ct, col_w, 0.34, fill=COLOR_CARD)
    txt(col_titles[i], cx+0.12, ct+0.04, col_w-0.2, 0.30,
        size=Pt(10.5), bold=True, color=COLOR_ACCENT2)

    # 本文
    txt_lines(col_bodies[i], cx+0.15, ct+0.42, col_w-0.25, 1.45,
              size=Pt(10.5), color=COLOR_LIGHT, spacing=Pt(3))

    # 結論バー
    rect(cx, ct+2.52, col_w, 0.36, fill=RGBColor(0x1A,0x1A,0x30))
    rect(cx, ct+2.52, 0.04, 0.36, fill=COLOR_ACCENT)
    txt(col_conclusions[i], cx+0.14, ct+2.56, col_w-0.2, 0.30,
        size=Pt(10), bold=True, color=COLOR_ACCENT)


# ============================================================
# 構造図：想定 vs 実態
# ============================================================
diag_t = 5.44
rect(0.5, diag_t, 12.5, 1.0,
     fill=RGBColor(0x18, 0x18, 0x2C),
     line=COLOR_DIVIDER, lw=Pt(0.5))

# ラベル
txt("営業の役割：想定 vs 実態", 0.65, diag_t+0.06, 3.5, 0.26,
    size=Pt(9.5), bold=True, color=COLOR_MUTED)

# 想定フロー
flow_t = diag_t + 0.36
boxes_ideal = [("営業アプローチ", COLOR_GREEN),
               ("案件創出", COLOR_GREEN),
               ("受注", COLOR_GREEN)]
boxes_real  = [("問い合わせ発生", COLOR_MUTED),
               ("営業が対応", COLOR_ACCENT2),
               ("受注", COLOR_ACCENT2)]

labels = [("想定（あるべき姿）", COLOR_GREEN, 0.62),
          ("実態（現実）",       COLOR_ACCENT, 5.92)]

for label, color, lx in labels:
    txt(label, lx, flow_t - 0.24, 2.5, 0.22,
        size=Pt(9), bold=True, color=color)

bw = 1.55
for j, (boxes, sx) in enumerate([(boxes_ideal, 0.62), (boxes_real, 5.92)]):
    for k, (label, color) in enumerate(boxes):
        bx = sx + k * (bw + 0.18)
        rect(bx, flow_t, bw, 0.40,
             fill=RGBColor(0x22,0x44,0x22) if color == COLOR_GREEN else RGBColor(0x33,0x28,0x18) if color == COLOR_ACCENT2 else RGBColor(0x28,0x28,0x40),
             line=color, lw=Pt(1.2))
        txt(label, bx+0.05, flow_t+0.06, bw-0.1, 0.30,
            size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)
        if k < 2:
            txt("→", bx+bw+0.01, flow_t+0.06, 0.20, 0.30,
                size=Pt(14), bold=True, color=color, align=PP_ALIGN.CENTER)

# VS バッジ
rect(5.32, flow_t, 0.55, 0.40, fill=COLOR_ACCENT)
txt("VS", 5.32, flow_t+0.05, 0.55, 0.32,
    size=Pt(12), bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# フッター：示唆
# ============================================================
rect(0, 6.55, 13.33, 0.95, fill=RGBColor(0x1A,0x0A,0x0A))
rect(0, 6.55, 13.33, 0.04, fill=COLOR_ACCENT)

txt("【示唆】",
    0.5, 6.62, 1.0, 0.30,
    size=Pt(10), bold=True, color=COLOR_ACCENT)

txt("現状の営業組織は「案件創出装置」ではなく「案件処理装置」として機能している。"
    "　再現性ある売上拡大のためには、営業の役割を「対応」から「創出」へ再定義し、"
    "起点データの記録・KPIの再設計・アカウント営業への転換が不可欠だ。",
    1.45, 6.62, 11.6, 0.80,
    size=Pt(10.5), color=COLOR_LIGHT)


# ============================================================
# 保存
# ============================================================
out = "/Volumes/SSD-PGU3C/営業推進/営業戦略スライド/BP事業部_現状営業分析.pptx"
prs.save(out)
print(f"Saved: {out}")
